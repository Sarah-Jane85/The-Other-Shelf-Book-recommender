"""
build_summary_slide.py
──────────────────────
Creates a single three-column summary slide matching the style of the
group Prezi: dark background, teal highlighted centre column.
Output: Non-Fiction Summary Slide.pptx (same folder)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours ────────────────────────────────────────────────────────────────
SLIDE_BG = RGBColor(0x1B, 0x2C, 0x2E)
COL_DARK = RGBColor(0x2D, 0x3D, 0x40)
COL_TEAL = RGBColor(0x1E, 0x95, 0x8B)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xCE, 0xE0, 0xDE)
SQ_DARK  = RGBColor(0x4A, 0x7E, 0x7C)
SQ_TEAL  = RGBColor(0x10, 0x58, 0x52)

HERE = Path(__file__).resolve().parent
OUT  = HERE / "Non-Fiction Summary Slide.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = SLIDE_BG

COL_W = 4.41
GAP   = 0.05
COLS  = [0, COL_W + GAP, 2 * (COL_W + GAP)]

# ── Column background rectangles ───────────────────────────────────────────
for l, color in zip(COLS, [COL_DARK, COL_TEAL, COL_DARK]):
    s = slide.shapes.add_shape(1, Inches(l), Inches(0), Inches(COL_W), Inches(7.5))
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = color


# ── Helpers ────────────────────────────────────────────────────────────────

def make_tf(col_idx, t=0.42, h=6.8):
    PAD = 0.22
    l   = COLS[col_idx] + PAD
    w   = COL_W - 2 * PAD
    tb  = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf  = tb.text_frame
    tf.word_wrap = True
    return tf


def add_run(para, text, size, bold=False, color=WHITE):
    r = para.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = color


def heading(tf, lines, size=27):
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        add_run(p, line, size, bold=True)


def spacer(tf, pts=7):
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = ""
    r.font.size = Pt(pts)


def bullet(tf, title, body, sq_color):
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    add_run(p, "▪  ", 12, color=sq_color)
    add_run(p, title, 16, bold=True)
    if body:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        add_run(p2, f"   {body}", 14, color=LIGHT)
    spacer(tf, 10)


# ══════════════════════════════════════════════════════════════════════════════
# COLUMN 1 — Genre Specific Challenges
# ══════════════════════════════════════════════════════════════════════════════
tf1 = make_tf(0)
heading(tf1, ["Genre Specific", "Challenges"], size=27)
spacer(tf1, 12)

bullet(tf1,
    "Language barrier:",
    "many core thinkers write in French & Portuguese "
    "— 1,211 books removed by language filter",
    SQ_DARK)

bullet(tf1,
    "Metadata ambiguity:",
    "memoir · history · biography · study guides "
    "· Wikipedia articles — hard to filter cleanly",
    SQ_DARK)

bullet(tf1,
    "Key challenge:",
    "defining left-wing non-fiction — the curated "
    "author list IS the genre boundary",
    SQ_DARK)


# ══════════════════════════════════════════════════════════════════════════════
# COLUMN 2 — Discovery  (highlighted)
# ══════════════════════════════════════════════════════════════════════════════
tf2 = make_tf(1)
heading(tf2, ["Discovery"], size=32)
spacer(tf2, 20)

bullet(tf2,
    "After cleaning, strong themes emerged:",
    "postcolonial theory · race · labour · feminist theory "
    "· decolonisation",
    SQ_TEAL)

bullet(tf2,
    "Insight:",
    "TF-IDF + cosine similarity works remarkably well "
    "in a focused, well-defined domain",
    SQ_TEAL)


# ══════════════════════════════════════════════════════════════════════════════
# COLUMN 3 — Takeaway
# ══════════════════════════════════════════════════════════════════════════════
tf3 = make_tf(2)
heading(tf3, ["Takeaway"], size=32)
spacer(tf3, 20)

bullet(tf3,
    "The recommender improved when the genre definition improved.",
    "",
    SQ_DARK)

bullet(tf3,
    "Data quality > model complexity",
    "",
    SQ_DARK)

bullet(tf3,
    "Bigrams captured compound concepts",
    '"racial capitalism" · "critical theory" — missed entirely by unigrams',
    SQ_DARK)


# ── Save ───────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved → {OUT}")
