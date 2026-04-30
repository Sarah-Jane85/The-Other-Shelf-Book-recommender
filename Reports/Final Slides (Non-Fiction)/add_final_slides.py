"""
add_final_slides.py
───────────────────
Opens the existing Non-Fiction Recommender.pptx, appends two new slides
(specific difficulties + learnings), and saves to this folder.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette (matches existing deck) ─────────────────────────────────
CREAM      = RGBColor(0xF5, 0xF0, 0xE8)
GOLD       = RGBColor(0xF5, 0xD7, 0x8E)
DARK_BROWN = RGBColor(0x2C, 0x1A, 0x0E)
MID_BROWN  = RGBColor(0x5C, 0x42, 0x08)
TAUPE      = RGBColor(0xD4, 0xC5, 0xA9)
DARK_BG    = RGBColor(0x1A, 0x10, 0x08)
CARD_BG    = RGBColor(0x2E, 0x1E, 0x10)

HERE   = Path(__file__).resolve().parent
SRC    = HERE.parent / "Slides (Non-Fiction)" / "Non-Fiction Recommender.pptx"
OUT    = HERE / "Non-Fiction Recommender.pptx"

prs = Presentation(SRC)
BLANK = prs.slide_layouts[6]
TOTAL = 10   # 8 existing + 2 new


# ── Helpers ────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill_color: RGBColor = None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def textbox(slide, text, l, t, w, h,
            size=24, bold=False, color=CREAM,
            align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def bullet_box(slide, items, l, t, w, h, size=18, color=CREAM):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = f"  •  {item}"
        run.font.size  = Pt(size)
        run.font.color.rgb = color
    return txBox


def divider(slide, t, color=GOLD):
    rect(slide, 0.4, t, 12.5, 0.04, fill_color=color)


def slide_number(slide, n):
    textbox(slide, f"{n} / {TOTAL}", 12.2, 7.1, 1.0, 0.35,
            size=11, color=TAUPE, align=PP_ALIGN.RIGHT)


def label(slide, text, l, t):
    textbox(slide, text.upper(), l, t, 6, 0.4,
            size=11, bold=True, color=GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Specific difficulties
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
rect(s, 0, 0, 0.35, 7.5, fill_color=MID_BROWN)
divider(s, 1.55)

label(s, "The Difficulties", 0.7, 0.25)
textbox(s, "What was specifically hard\nabout this genre",
        0.7, 0.55, 11, 1.3,
        size=36, bold=True, color=CREAM)

# Left column
textbox(s, "Genre & data", 0.7, 1.85, 5.8, 0.45,
        size=15, bold=True, color=GOLD)
bullet_box(s, [
    "Language barrier: thinkers like Mbembe, Fanon, and Césaire write "
    "primarily in French or Portuguese — langdetect removed 1,211 books, often the most important ones.",
    "Some key authors returned zero API results (W.E.B. Du Bois, Saidiya Hartman, "
    "C.L.R. James) — likely due to name-format mismatches in Open Library.",
    "The subjects column was 100% null in Goodreads data — no topic tags to "
    "use as a fallback signal.",
    "Curating 39 authors is itself a subjective decision: who gets included "
    "shapes what the recommender can and cannot find.",
], 0.7, 2.3, 5.8, 4.3, size=15, color=CREAM)

# Right column
textbox(s, "Pipeline & engineering", 6.8, 1.85, 6.0, 0.45,
        size=15, bold=True, color=GOLD)
bullet_box(s, [
    "Description enrichment: 1,314 of 5,487 books had no synopsis — each required "
    "an individual Goodreads page visit, with progress saved every 50 books to "
    "survive crashes.",
    "Merging two sources with incompatible key formats: Open Library uses "
    "/works/OL… identifiers; Goodreads uses /book/show/… paths — cross-linking "
    "them required careful parsing.",
    "No ground truth: TF-IDF is unsupervised and quality evaluation was entirely "
    "manual — hard to know whether tuning ngram range or max_features actually improved results.",
], 6.8, 2.3, 6.0, 4.3, size=15, color=CREAM)

rect(s, 6.5, 1.8, 0.04, 5.1, fill_color=MID_BROWN)
slide_number(s, 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Learnings
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
rect(s, 0, 0, 0.35, 7.5, fill_color=MID_BROWN)
divider(s, 1.55)

label(s, "The Learnings", 0.7, 0.25)
textbox(s, "What this project taught me",
        0.7, 0.55, 11, 1.0,
        size=36, bold=True, color=CREAM)

# Four cards in a 2×2 grid
cards = [
    ("Vectorisation matters",
     "Including title + author + description in TF-IDF outperformed title alone "
     "significantly. Bigrams (ngram_range=(1,2)) captured compound terms like "
     "\"racial capitalism\" and \"critical theory\" that unigrams miss entirely."),
    ("Scraping needs resilience",
     "At scale, crashes and rate limits are inevitable. Saving progress every "
     "50 books and throttling requests turned a fragile script into a reliable "
     "pipeline. Patience is a data engineering skill."),
    ("Curation is design",
     "Choosing which 39 authors to include shapes what the recommender can "
     "surface as much as the algorithm does. Exclusions are invisible — but they "
     "define the edges of the system."),
    ("Simple works in a focused domain",
     "TF-IDF + cosine similarity is decades-old technology, but it produces "
     "genuinely useful recommendations when the data is clean and the domain is "
     "well-scoped. Complexity isn't always the answer."),
]

card_w = 5.9
for i, (title, body) in enumerate(cards):
    row = i // 2
    col = i % 2
    l = 0.7  + col * 6.3
    t = 2.15 + row * 2.5
    rect(s, l, t, card_w, 2.2, fill_color=CARD_BG)
    textbox(s, title, l + 0.15, t + 0.12, card_w - 0.25, 0.45,
            size=14, bold=True, color=GOLD)
    textbox(s, body,  l + 0.15, t + 0.58, card_w - 0.25, 1.5,
            size=13, color=CREAM)

slide_number(s, 10)


# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved → {OUT}")
