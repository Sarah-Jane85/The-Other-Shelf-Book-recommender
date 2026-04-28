"""
build_slides.py
───────────────
Generates the Non-Fiction presentation as a .pptx file.
Run from anywhere — output lands in the same folder as this script.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ─────────────────────────────────────────────────────────
CREAM      = RGBColor(0xF5, 0xF0, 0xE8)
GOLD       = RGBColor(0xF5, 0xD7, 0x8E)
DARK_BROWN = RGBColor(0x2C, 0x1A, 0x0E)
MID_BROWN  = RGBColor(0x5C, 0x42, 0x08)
TAUPE      = RGBColor(0xD4, 0xC5, 0xA9)
DARK_BG    = RGBColor(0x1A, 0x10, 0x08)

HERE = Path(__file__).resolve().parent
OUT  = HERE / "Non-Fiction Recommender.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helpers ────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill_color: RGBColor = None, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def textbox(slide, text, l, t, w, h,
            size=24, bold=False, color=CREAM,
            align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf   = txBox.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run  = para.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def bullet_box(slide, items, l, t, w, h,
               size=20, color=CREAM, indent=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        bullet = "      •  " if indent else "  •  "
        run = para.add_run()
        run.text = f"{bullet}{item}"
        run.font.size  = Pt(size)
        run.font.color.rgb = color

    return txBox


def divider(slide, t, color=GOLD):
    rect(slide, 0.4, t, 12.5, 0.04, fill_color=color)


def slide_number(slide, n, total=7):
    textbox(slide, f"{n} / {total}", 12.2, 7.1, 1.0, 0.35,
            size=11, color=TAUPE, align=PP_ALIGN.RIGHT)


def label(slide, text, l, t):
    textbox(slide, text.upper(), l, t, 4, 0.4,
            size=11, bold=True, color=GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Why this genre?
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
rect(s, 0, 0, 0.35, 7.5, fill_color=MID_BROWN)
divider(s, 1.55)

label(s, "The Genre", 0.7, 0.25)
textbox(s, "Why Non-Fiction — and why\nthis corner of it?",
        0.7, 0.6, 11, 1.3,
        size=38, bold=True, color=CREAM)

bullet_box(s, [
    "Non-fiction is dominated by self-help, business, and biographies — "
    "that's what algorithms push.",
    "Left-wing thought, critical theory, and postcolonial studies form a rich "
    "intellectual tradition that gets almost no mainstream visibility.",
    "Navigating this niche is genuinely hard: overlapping topics, thousands of authors, "
    "no clear entry points.",
    "Capitalism, decolonisation, race, labour, feminist theory — these conversations matter, "
    "but the books are hard to find.",
    "I wanted to build something that treats this genre seriously, not as a footnote.",
], 0.7, 2.0, 12.0, 4.8, size=19, color=CREAM)

slide_number(s, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Why a recommender?
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
rect(s, 0, 0, 0.35, 7.5, fill_color=MID_BROWN)
divider(s, 1.55)

label(s, "The Problem", 0.7, 0.25)
textbox(s, "Why does a recommender help?",
        0.7, 0.6, 11, 1.0,
        size=38, bold=True, color=CREAM)

bullet_box(s, [
    "Nobody has time to browse 3,000 books.",
    "Goodreads doesn't know this niche.",
    "Amazon optimises for bestsellers, not theory.",
    "Google gives you Wikipedia, not the next book to read.",
], 0.7, 2.05, 6.0, 3.5, size=19, color=CREAM)

bullet_box(s, [
    "New to the genre? Hard to know where to start.",
    "Already deep in it? Hard to know what comes next.",
    "Topic overlap is massive — postcolonial theory, Marxism, "
    "feminist theory, critical race theory all intersect.",
    "A recommender connects the dots for you.",
], 6.9, 2.05, 6.0, 3.5, size=19, color=CREAM)

rect(s, 6.6, 2.0, 0.04, 4.8, fill_color=MID_BROWN)

textbox(s,
        "→  The recommender saves time, surfaces the unexpected, and helps readers "
        "navigate a dense and rewarding genre.",
        0.7, 5.9, 12.0, 0.9,
        size=17, color=GOLD, italic=True)

slide_number(s, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — How it works
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
rect(s, 0, 0, 0.35, 7.5, fill_color=MID_BROWN)
divider(s, 1.55)

label(s, "The Method", 0.7, 0.25)
textbox(s, "How the recommender works",
        0.7, 0.6, 11, 1.0,
        size=38, bold=True, color=CREAM)

steps = [
    ("01", "Collect",    "~3,034 books scraped from APIs and the web.\nTitles, authors, descriptions, publication years."),
    ("02", "Clean",      "Deduplication, missing-value handling,\nnormalising formats across sources."),
    ("03", "Vectorise",  "TF-IDF on titles + descriptions.\nConverts text into numerical vectors."),
    ("04", "Match",      "Cosine similarity between query vector\nand every book in the index."),
    ("05", "Serve",      "Top-N results ranked by similarity,\ndeployed as a Streamlit web app."),
]

box_w = 2.3
for i, (num, title, body) in enumerate(steps):
    l = 0.55 + i * (box_w + 0.12)
    rect(s, l, 2.1, box_w, 4.5, fill_color=RGBColor(0x2E, 0x1E, 0x10))
    textbox(s, num,   l + 0.12, 2.2,  box_w - 0.2, 0.55, size=28, bold=True, color=GOLD)
    textbox(s, title, l + 0.12, 2.75, box_w - 0.2, 0.55, size=17, bold=True, color=CREAM)
    textbox(s, body,  l + 0.12, 3.35, box_w - 0.2, 2.8,  size=14, color=TAUPE)

slide_number(s, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Methodology
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
rect(s, 0, 0, 0.35, 7.5, fill_color=MID_BROWN)
divider(s, 1.55)

label(s, "The Methodology", 0.7, 0.25)
textbox(s, "Why content-based filtering?",
        0.7, 0.6, 11, 1.0,
        size=38, bold=True, color=CREAM)

textbox(s, "The approach", 0.7, 1.75, 5.8, 0.5,
        size=17, bold=True, color=GOLD)
bullet_box(s, [
    "Content-based filtering — recommends based on what a book is about, "
    "not on what other users clicked.",
    "Each book is turned into a vector of weighted terms (TF-IDF). "
    "The query is vectorised the same way.",
    "Cosine similarity measures the angle between vectors — it captures "
    "topical overlap regardless of book length.",
    "Returns the top-N closest books in the vector space.",
], 0.7, 2.2, 5.8, 4.0, size=17, color=CREAM)

textbox(s, "Why not collaborative filtering?", 6.8, 1.75, 6.0, 0.5,
        size=17, bold=True, color=GOLD)
bullet_box(s, [
    "Collaborative filtering needs user ratings and interaction history — "
    "we have neither.",
    "This niche is too small for meaningful user-behaviour data.",
    "TF-IDF is interpretable: you can see exactly which terms drove a recommendation.",
    "EDA used KMeans clustering (k=5) + PCA to validate that the books naturally "
    "group into coherent topic clusters — confirming the vector space is meaningful.",
], 6.8, 2.2, 6.0, 4.0, size=17, color=CREAM)

rect(s, 6.5, 1.7, 0.04, 5.0, fill_color=MID_BROWN)
slide_number(s, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Our process
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
rect(s, 0, 0, 0.35, 7.5, fill_color=MID_BROWN)
divider(s, 1.55)

label(s, "The Journey", 0.7, 0.25)
textbox(s, "What the process actually looked like",
        0.7, 0.6, 11, 1.0,
        size=38, bold=True, color=CREAM)

phases = [
    ("Data collection",
     ["Hit API rate limits repeatedly — had to slow down and batch requests.",
      "Scraped two sources and merged them, with lots of duplicate noise."]),
    ("Data cleaning",
     ["Inconsistent author names, missing descriptions, broken years.",
      "Manually reviewed edge cases — some 'books' were study guides or Wikipedia dumps."]),
    ("Building the model",
     ["TF-IDF is simple but surprisingly effective for text matching.",
      "Tuning what to include in the vectoriser (title only vs title + description) made a big difference."]),
    ("The app",
     ["Connecting the recommender to Streamlit was smooth.",
      "Cover images from Open Library added life to the results — but not all books have covers."]),
]

for i, (phase, points) in enumerate(phases):
    row = i // 2
    col = i %  2
    l = 0.7  + col * 6.3
    t = 2.15 + row * 2.4
    rect(s, l, t, 5.9, 2.1, fill_color=RGBColor(0x2E, 0x1E, 0x10))
    textbox(s, phase, l + 0.15, t + 0.12, 5.6, 0.45, size=15, bold=True, color=GOLD)
    bullet_box(s, points, l + 0.05, t + 0.55, 5.75, 1.4, size=13, color=CREAM)

slide_number(s, 5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — The product & audience
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
rect(s, 0, 0, 0.35, 7.5, fill_color=MID_BROWN)
divider(s, 1.55)

label(s, "The Product", 0.7, 0.25)
textbox(s, "Who it's for — and why it matters",
        0.7, 0.6, 11, 1.0,
        size=38, bold=True, color=CREAM)

textbox(s, "The audience", 0.7, 1.75, 5.8, 0.5,
        size=17, bold=True, color=GOLD)
bullet_box(s, [
    "Readers who already know they want 'something like Naomi Klein'.",
    "Students and researchers navigating critical theory for the first time.",
    "Activists and organisers looking for reading lists.",
    "Anyone who's frustrated that Goodreads keeps recommending the same 10 books.",
], 0.7, 2.2, 5.8, 3.8, size=17, color=CREAM)

textbox(s, "Why it's valuable", 6.8, 1.75, 6.0, 0.5,
        size=17, bold=True, color=GOLD)
bullet_box(s, [
    "Saves hours of browsing and searching.",
    "Surfaces authors you'd never find on your own.",
    "Works by topic, author, or free-text keyword — flexible for different needs.",
    "Brings a marginalised genre into the same UX quality as mainstream recommenders.",
    "Free, open, and accessible in the browser — no account needed.",
], 6.8, 2.2, 6.0, 3.8, size=17, color=CREAM)

rect(s, 6.5, 1.7, 0.04, 5.0, fill_color=MID_BROWN)
slide_number(s, 6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Challenges
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
rect(s, 0, 0, 0.35, 7.5, fill_color=MID_BROWN)
divider(s, 1.55)

label(s, "The Challenges", 0.7, 0.25)
textbox(s, "What was hard — and what we learned",
        0.7, 0.6, 11, 1.0,
        size=38, bold=True, color=CREAM)

challenges = [
    ("Scraping at scale",
     ["APIs imposed rate limits — every source had to be batched and throttled.",
      "Scraping ~3,000 books took significantly longer than expected.",
      "Two sources had different formats; merging them cleanly required careful alignment."]),
    ("Dirty data",
     ["Author names had dozens of variants for the same person.",
      "Descriptions were missing for a large portion of titles.",
      "Some entries weren't books at all — study guides, Wikipedia articles, duplicates."]),
    ("Choosing what to vectorise",
     ["Title-only vectorisation was too shallow; descriptions added crucial context.",
      "Stop-word removal and n-gram tuning had a noticeable impact on result quality.",
      "No ground-truth labels — evaluation was manual and subjective."]),
    ("Cover images & links",
     ["Not all books exist in Open Library — missing covers needed a graceful fallback.",
      "Goodreads links depended on Open Library keys being in the right format.",
      "Live API calls in Streamlit required caching to avoid slow load times."]),
]

for i, (challenge, points) in enumerate(challenges):
    row = i // 2
    col = i %  2
    l = 0.7  + col * 6.3
    t = 2.15 + row * 2.35
    rect(s, l, t, 5.9, 2.1, fill_color=RGBColor(0x2E, 0x1E, 0x10))
    textbox(s, challenge, l + 0.15, t + 0.12, 5.6, 0.45, size=15, bold=True, color=GOLD)
    bullet_box(s, points, l + 0.05, t + 0.55, 5.75, 1.45, size=12, color=CREAM)

slide_number(s, 7)


# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved → {OUT}")
